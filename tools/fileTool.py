import os
import json
import subprocess
import platform
from pathlib import Path
from typing import Optional, List, Dict, Any
from agno.tools import Toolkit
from agno.tools.decorator import tool

# 文档解析相关库
try:
    from docx import Document  # 解析 .docx

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl  # 解析 .xlsx

    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation  # 解析 .pptx

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import olefile  # 解析旧版 .doc 和 .wps
    import textract  # 备用解析方案

    OLE_AVAILABLE = True
except ImportError:
    OLE_AVAILABLE = False

try:
    from PIL import Image  # 解析图片
    import pytesseract  # OCR 识别图片中的文字

    OCR_AVAILABLE = False  # 默认关闭，需要安装 tesseract
except ImportError:
    OCR_AVAILABLE = False


class FileTool(Toolkit):
    """文件操作工具集 - 让 Agent 读取、打开、管理本地文件（支持 WPS 多种格式）"""

    def __init__(self, base_path: Optional[str] = None, max_file_size_mb: int = 50,
                 enable_ocr: bool = False):
        """
        初始化文件工具

        Args:
            base_path: 基础路径限制（安全考虑），None 表示允许访问任意路径
            max_file_size_mb: 最大允许读取的文件大小（MB），WPS文档可能较大，默认50MB
            enable_ocr: 是否启用图片OCR文字识别（需要安装tesseract）
        """
        super().__init__(name="file_tool")
        self.base_path = Path(base_path) if base_path else None
        self.max_file_size = max_file_size_mb * 1024 * 1024
        self.enable_ocr = enable_ocr and OCR_AVAILABLE

        # 检查依赖库状态
        self.lib_status = {
            "docx": DOCX_AVAILABLE,
            "xlsx": XLSX_AVAILABLE,
            "pptx": PPTX_AVAILABLE,
            "ole": OLE_AVAILABLE,
            "ocr": self.enable_ocr
        }

        # 注册所有工具方法
        self.register(self.read_file)
        self.register(self.list_directory)
        self.register(self.get_file_info)
        self.register(self.search_in_file)
        self.register(self.read_multiple_files)
        self.register(self.open_file)
        self.register(self.open_folder)
        self.register(self.check_format_support)  # 新增：检查格式支持

    def _validate_path(self, file_path: str) -> Path:
        """验证路径安全性"""
        path = Path(file_path).expanduser().resolve()

        # 检查文件大小限制（仅对文件有效）
        if path.exists() and path.is_file():
            file_size = path.stat().st_size
            if file_size > self.max_file_size:
                raise ValueError(
                    f"文件过大（{file_size / 1024 / 1024:.1f}MB），超过限制 {self.max_file_size / 1024 / 1024:.0f}MB")

        # 检查基础路径限制
        if self.base_path:
            try:
                path.relative_to(self.base_path)
            except ValueError:
                raise PermissionError(f"禁止访问：{path} 不在允许的目录 {self.base_path} 内")

        return path

    def _get_supported_extensions(self) -> Dict[str, str]:
        """获取支持的文件扩展名和对应的解析方法"""
        extensions = {
            # 文本文件
            '.txt': 'text', '.md': 'text', '.csv': 'text', '.json': 'json',
            '.py': 'code', '.js': 'code', '.html': 'code', '.css': 'code',
            '.xml': 'text', '.log': 'text', '.ini': 'text', '.cfg': 'text',

            # WPS/Microsoft Office 文档
            '.docx': 'docx', '.doc': 'doc_old', '.wps': 'wps',
            '.xlsx': 'xlsx', '.xls': 'xls_old',
            '.pptx': 'pptx', '.ppt': 'ppt_old',

            # PDF
            '.pdf': 'pdf',

            # 图片（可选OCR）
            '.jpg': 'image', '.jpeg': 'image', '.png': 'image',
            '.bmp': 'image', '.tiff': 'image', '.gif': 'image',
        }
        return extensions

    @tool
    def check_format_support(self, file_path: str) -> str:
        """
        检查文件格式是否支持解析，并返回支持情况

        Args:
            file_path: 文件路径

        Returns:
            格式支持情况说明
        """
        try:
            path = self._validate_path(file_path)
            suffix = path.suffix.lower()
            extensions = self._get_supported_extensions()

            if suffix not in extensions:
                return f"⚠️ 格式 {suffix} 暂不支持解析，但可以尝试用系统程序打开"

            file_type = extensions[suffix]

            # 检查是否需要特定库
            required_libs = {
                'docx': ('python-docx', DOCX_AVAILABLE),
                'xlsx': ('openpyxl', XLSX_AVAILABLE),
                'pptx': ('python-pptx', PPTX_AVAILABLE),
                'doc_old': ('olefile + textract', OLE_AVAILABLE),
                'wps': ('olefile + textract', OLE_AVAILABLE),
                'xls_old': ('olefile + textract', OLE_AVAILABLE),
                'ppt_old': ('olefile + textract', OLE_AVAILABLE),
                'image': ('PIL + pytesseract', self.enable_ocr),
            }

            if file_type in required_libs:
                lib_name, is_available = required_libs[file_type]
                if not is_available:
                    return f"📄 格式 {suffix} 需要安装 {lib_name} 才能解析文本内容，但可以用 open_file 工具打开查看"

            return f"✅ 格式 {suffix} 支持解析"

        except Exception as e:
            return f"检查失败：{str(e)}"

    def _parse_docx(self, file_path: Path) -> str:
        """解析 .docx 文件"""
        if not DOCX_AVAILABLE:
            return None, "需要安装 python-docx: pip install python-docx"

        try:
            doc = Document(str(file_path))
            content = []

            # 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)

            # 提取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content.append(" | ".join(row_text))

            return "\n".join(content), None
        except Exception as e:
            return None, f"解析失败：{str(e)}"

    def _parse_xlsx(self, file_path: Path) -> str:
        """解析 .xlsx 文件"""
        if not XLSX_AVAILABLE:
            return None, "需要安装 openpyxl: pip install openpyxl"

        try:
            workbook = openpyxl.load_workbook(str(file_path), data_only=True)
            content = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content.append(f"\n📊 工作表：{sheet_name}")
                content.append("-" * 50)

                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # 只显示非空行
                        content.append(" | ".join(row_data))

            return "\n".join(content), None
        except Exception as e:
            return None, f"解析失败：{str(e)}"

    def _parse_pptx(self, file_path: Path) -> str:
        """解析 .pptx 文件"""
        if not PPTX_AVAILABLE:
            return None, "需要安装 python-pptx: pip install python-pptx"

        try:
            prs = Presentation(str(file_path))
            content = []

            for idx, slide in enumerate(prs.slides, 1):
                content.append(f"\n📽️ 幻灯片 {idx}")
                content.append("-" * 50)

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text)

            return "\n".join(content), None
        except Exception as e:
            return None, f"解析失败：{str(e)}"

    def _parse_old_doc(self, file_path: Path, file_type: str) -> str:
        """解析旧版 .doc, .wps, .xls, .ppt 文件（使用 textract 或 olefile）"""
        if not OLE_AVAILABLE:
            return None, "需要安装 textract: pip install textract 或 pip install olefile"

        try:
            # 尝试使用 textract（支持多种格式）
            import textract
            text = textract.process(str(file_path)).decode('utf-8', errors='ignore')
            return text, None
        except ImportError:
            # 降级方案：只读取基本信息
            try:
                import olefile
                if olefile.isOleFile(str(file_path)):
                    ole = olefile.OleFileIO(str(file_path))
                    # 获取文件基本信息
                    info = [f"📄 {file_type.upper()} 文件（旧版格式）"]
                    info.append(f"文件大小：{self._format_size(file_path.stat().st_size)}")
                    info.append("\n⚠️ 需要安装 textract 才能提取文本内容：")
                    info.append("  pip install textract")
                    info.append("\n或者将文件另存为 .docx/.xlsx/.pptx 格式")
                    return "\n".join(info), None
                return None, "不是有效的 OLE 文件"
            except:
                return None, "无法解析旧版文档格式"

    def _parse_image(self, file_path: Path) -> str:
        """解析图片文件（提取文字）"""
        if not self.enable_ocr:
            return None, "OCR 未启用，如需识别图片文字请设置 enable_ocr=True 并安装 pytesseract"

        try:
            from PIL import Image
            import pytesseract

            image = Image.open(str(file_path))
            # 使用中文+英文语言包
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')

            if text.strip():
                return f"🖼️ 图片中的文字识别结果：\n\n{text}", None
            else:
                return None, "图片中未检测到文字"
        except Exception as e:
            return None, f"OCR 识别失败：{str(e)}"

    def _parse_pdf(self, file_path: Path) -> str:
        """解析 PDF 文件（需要 PyPDF2 或 pdfplumber）"""
        try:
            # 尝试使用 PyPDF2
            import PyPDF2
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                content = []
                for page_num, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        content.append(f"--- 第 {page_num} 页 ---")
                        content.append(text)
                return "\n".join(content), None
        except ImportError:
            try:
                # 尝试使用 pdfplumber（更准确）
                import pdfplumber
                with pdfplumber.open(str(file_path)) as pdf:
                    content = []
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text:
                            content.append(f"--- 第 {page_num} 页 ---")
                            content.append(text)
                    return "\n".join(content), None
            except ImportError:
                return None, "需要安装 PyPDF2 或 pdfplumber: pip install PyPDF2"

    @tool
    def read_file(self, file_path: str, encoding: str = "utf-8", max_lines: Optional[int] = None) -> str:
        """
        读取文件内容（支持多种格式：TXT、DOCX、XLSX、PPTX、WPS、PDF、图片等）

        Args:
            file_path: 文件路径（绝对路径或相对路径）
            encoding: 文件编码，默认 utf-8（仅对文本文件有效）
            max_lines: 最大读取行数，None 表示读取全部（仅对文本文件有效）

        Returns:
            文件内容字符串
        """
        try:
            path = self._validate_path(file_path)

            if not path.exists():
                return f"❌ 文件不存在：{file_path}"

            if not path.is_file():
                return f"❌ 路径不是文件：{file_path}"

            suffix = path.suffix.lower()
            extensions = self._get_supported_extensions()

            # 获取文件类型
            file_type = extensions.get(suffix, 'unknown')

            file_size = path.stat().st_size
            size_str = self._format_size(file_size)

            # 根据文件类型选择解析方法
            content = None
            warning = None

            # 文本文件
            if file_type in ['text', 'code']:
                try:
                    with open(path, 'r', encoding=encoding, errors='ignore') as f:
                        if max_lines:
                            lines = [next(f) for _ in range(max_lines)]
                            content = ''.join(lines)
                            if len(lines) == max_lines:
                                content += f"\n... (文件内容超过 {max_lines} 行，已截断)"
                        else:
                            content = f.read()
                except UnicodeDecodeError:
                    for enc in ['gbk', 'gb2312', 'latin-1']:
                        try:
                            with open(path, 'r', encoding=enc) as f:
                                content = f.read()
                            warning = f"使用 {enc} 编码"
                            break
                        except UnicodeDecodeError:
                            continue
                    if content is None:
                        return f"❌ 无法解码文件：{file_path}"

            # JSON 文件
            elif file_type == 'json':
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    data = json.load(f)
                    content = json.dumps(data, ensure_ascii=False, indent=2)

            # Word 文档
            elif file_type == 'docx':
                content, error = self._parse_docx(path)
                if error:
                    return f"❌ {error}\n\n💡 提示：可以用 open_file 工具用 Word/WPS 打开查看"

            # Excel 文档
            elif file_type == 'xlsx':
                content, error = self._parse_xlsx(path)
                if error:
                    return f"❌ {error}\n\n💡 提示：可以用 open_file 工具用 Excel/WPS 打开查看"

            # PowerPoint 文档
            elif file_type == 'pptx':
                content, error = self._parse_pptx(path)
                if error:
                    return f"❌ {error}\n\n💡 提示：可以用 open_file 工具用 PowerPoint/WPS 打开查看"

            # 旧版 Office 文档
            elif file_type in ['doc_old', 'wps', 'xls_old', 'ppt_old']:
                content, error = self._parse_old_doc(path, file_type)
                if error:
                    return f"❌ {error}\n\n💡 提示：可以用 open_file 工具用 WPS 打开查看"

            # PDF 文档
            elif file_type == 'pdf':
                content, error = self._parse_pdf(path)
                if error:
                    return f"❌ {error}\n\n💡 提示：可以用 open_file 工具用 PDF 阅读器打开查看"

            # 图片文件
            elif file_type == 'image':
                content, error = self._parse_image(path)
                if error:
                    return f"❌ {error}\n\n💡 提示：可以用 open_file 工具用图片查看器打开"

            else:
                return f"📄 文件格式 {suffix} 暂不支持解析文本内容\n\n" \
                       f"文件大小：{size_str}\n" \
                       f"文件路径：{path.absolute()}\n\n" \
                       f"💡 提示：\n" \
                       f"  1. 使用 open_file 工具用默认程序打开查看\n" \
                       f"  2. 如果文件是文本格式，可以尝试指定其他编码\n" \
                       f"  3. 或安装相应解析库后获得支持"

            # 格式化输出
            result = f"📄 文件：{path.name}\n📏 大小：{size_str}\n"
            if warning:
                result += f"⚠️ 提示：{warning}\n"
            result += f"\n📝 内容：\n{'-' * 50}\n{content}"

            # 对于长内容进行截断提示
            if len(result) > 10000:
                result = result[:10000] + "\n\n... (内容过长，已截断，建议使用 search_in_file 搜索关键词)"

            return result

        except PermissionError as e:
            return f"❌ 权限不足：{str(e)}"
        except Exception as e:
            return f"❌ 读取失败：{str(e)}"

    @tool
    def list_directory(self, directory_path: str = ".", show_hidden: bool = False) -> str:
        """列出目录中的文件和文件夹"""
        try:
            path = self._validate_path(directory_path)

            if not path.exists():
                return f"❌ 目录不存在：{directory_path}"

            if not path.is_dir():
                return f"❌ 路径不是目录：{directory_path}"

            items = []
            for item in path.iterdir():
                if not show_hidden and item.name.startswith('.'):
                    continue

                item_type = "📁" if item.is_dir() else "📄"
                size = ""
                if item.is_file():
                    size = f" ({self._format_size(item.stat().st_size)})"
                items.append(f"{item_type} {item.name}{size}")

            if not items:
                return f"📁 目录为空：{directory_path}"

            return f"📁 目录：{path.absolute()}\n\n" + "\n".join(sorted(items))

        except Exception as e:
            return f"❌ 列表失败：{str(e)}"

    @tool
    def get_file_info(self, file_path: str) -> str:
        """获取文件的详细信息（大小、修改时间、类型等）"""
        try:
            path = self._validate_path(file_path)

            if not path.exists():
                return f"❌ 文件不存在：{file_path}"

            stat = path.stat()
            suffix = path.suffix.lower()
            extensions = self._get_supported_extensions()

            info = {
                "文件名": path.name,
                "绝对路径": str(path.absolute()),
                "类型": "📁 目录" if path.is_dir() else "📄 文件",
                "格式支持": extensions.get(suffix, '暂不支持'),
                "大小": self._format_size(stat.st_size),
                "创建时间": self._format_time(stat.st_ctime),
                "修改时间": self._format_time(stat.st_mtime),
                "访问时间": self._format_time(stat.st_atime),
            }

            if path.is_file():
                info["扩展名"] = path.suffix or "无"
                info["是否可读"] = "是" if os.access(path, os.R_OK) else "否"
                info["是否可写"] = "是" if os.access(path, os.W_OK) else "否"

            result = "📋 文件信息：\n\n"
            for key, value in info.items():
                result += f"  {key}：{value}\n"

            return result

        except Exception as e:
            return f"❌ 获取信息失败：{str(e)}"

    @tool
    def search_in_file(self, file_path: str, keyword: str, context_lines: int = 2) -> str:
        """在文件中搜索关键词（支持文本文件）"""
        try:
            path = self._validate_path(file_path)

            if not path.exists():
                return f"❌ 文件不存在：{file_path}"

            # 只对文本文件进行搜索
            suffix = path.suffix.lower()
            if suffix not in ['.txt', '.md', '.py', '.json', '.csv', '.log', '.html', '.css', '.js']:
                return f"⚠️ 暂不支持在 {suffix} 格式文件中搜索，建议先使用 read_file 读取内容"

            content = path.read_text(encoding='utf-8', errors='ignore')
            lines = content.split('\n')

            matches = []
            for i, line in enumerate(lines, 1):
                if keyword.lower() in line.lower():
                    matches.append({
                        "line_num": i,
                        "line": line.strip()[:200]
                    })

            if not matches:
                return f"🔍 未找到关键词 '{keyword}'"

            result = f"🔍 在 {path.name} 中找到 {len(matches)} 处匹配 '{keyword}'：\n\n"
            for match in matches[:20]:  # 最多显示20处
                result += f"  行 {match['line_num']}: {match['line']}\n"

            if len(matches) > 20:
                result += f"\n  ... 还有 {len(matches) - 20} 处匹配"

            return result

        except Exception as e:
            return f"❌ 搜索失败：{str(e)}"

    @tool
    def read_multiple_files(self, file_paths: str, encoding: str = "utf-8") -> str:
        """同时读取多个文件（用逗号分隔路径）"""
        paths = [p.strip() for p in file_paths.split(',')]
        results = []

        for file_path in paths:
            result = self.read_file(file_path, encoding)
            results.append(result)

        if not results:
            return "❌ 没有成功读取任何文件"

        return "\n" + "=" * 60 + "\n".join(results)

    @tool
    def open_file(self, file_path: str) -> str:
        """用系统默认应用程序打开文件（适合查看 WPS 文档）"""
        try:
            path = self._validate_path(file_path)

            if not path.exists():
                return f"❌ 文件不存在：{file_path}"

            if not path.is_file():
                return f"❌ 路径不是文件：{file_path}"

            system = platform.system()

            if system == "Windows":
                os.startfile(str(path))
            elif system == "Darwin":  # macOS
                subprocess.run(["open", str(path)], check=True)
            elif system == "Linux":
                subprocess.run(["xdg-open", str(path)], check=True)
            else:
                return f"❌ 不支持的操作系统：{system}"

            return f"✅ 已用默认程序打开文件：{path.name}\n💡 提示：如果文件未正常显示，请确保已安装 WPS Office 或 Microsoft Office"

        except subprocess.CalledProcessError as e:
            return f"❌ 打开失败：{str(e)}\n💡 提示：请检查是否已安装相应的办公软件"
        except PermissionError as e:
            return f"❌ 权限不足：{str(e)}"
        except Exception as e:
            return f"❌ 打开失败：{str(e)}"

    @tool
    def open_folder(self, folder_path: str = ".") -> str:
        """用系统文件管理器打开文件夹"""
        try:
            path = self._validate_path(folder_path)

            if not path.exists():
                return f"❌ 目录不存在：{folder_path}"

            if not path.is_dir():
                return f"❌ 路径不是目录：{folder_path}"

            system = platform.system()

            if system == "Windows":
                os.startfile(str(path))
            elif system == "Darwin":  # macOS
                subprocess.run(["open", str(path)], check=True)
            elif system == "Linux":
                subprocess.run(["xdg-open", str(path)], check=True)
            else:
                return f"❌ 不支持的操作系统：{system}"

            return f"✅ 已打开文件夹：{path.absolute()}"

        except Exception as e:
            return f"❌ 打开失败：{str(e)}"

    def _format_size(self, size_bytes: int) -> str:
        """格式化文件大小"""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size_bytes < 1024:
                return f"{size_bytes:.1f} {unit}"
            size_bytes /= 1024
        return f"{size_bytes:.1f} TB"

    def _format_time(self, timestamp: float) -> str:
        """格式化时间戳"""
        from datetime import datetime
        return datetime.fromtimestamp(timestamp).strftime("%Y-%m-%d %H:%M:%S")


# 依赖安装说明
INSTALL_GUIDE = """
📦 安装依赖以支持更多格式：

# 基础文本文件（无需额外安装）

# Word .docx 文件
pip install python-docx

# Excel .xlsx 文件  
pip install openpyxl

# PowerPoint .pptx 文件
pip install python-pptx

# 旧版 Office 文件 (.doc, .wps, .xls, .ppt)
pip install textract  # Windows可能需要安装其他依赖

# PDF 文件
pip install PyPDF2
# 或更精确的解析
pip install pdfplumber

# 图片 OCR（识别图片中的文字）
pip install pillow pytesseract
# 还需要安装 tesseract-ocr 引擎
# Windows: https://github.com/UB-Mannheim/tesseract/wiki
# macOS: brew install tesseract tesseract-lang
# Linux: sudo apt-get install tesseract-ocr tesseract-ocr-chi-sim

# 或者一次性安装所有依赖
pip install python-docx openpyxl python-pptx textract PyPDF2 pdfplumber pillow pytesseract
"""